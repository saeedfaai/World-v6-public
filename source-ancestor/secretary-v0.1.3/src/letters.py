from __future__ import annotations
from pathlib import Path
import shutil, subprocess, tempfile
from pptx import Presentation
from pptx.util import Inches, Pt


def _replace_single_run(shape, text: str) -> None:
    p = shape.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.text = text


def _replace_body(shape, body: str) -> None:
    ps = shape.text_frame.paragraphs
    values = ["بسمه تعالی", "", "با سلام و احترام،", body, "", "", "", ""]
    for i, value in enumerate(values):
        if i >= len(ps):
            break
        if ps[i].runs:
            ps[i].runs[0].text = value
            for r in ps[i].runs[1:]:
                r.text = ""
        else:
            ps[i].text = value


def _replace_signature(shape, signatory: str, company: str) -> None:
    ps = shape.text_frame.paragraphs
    vals = ["با تشکر و احترام", signatory, company]
    for p, value in zip(ps, vals):
        if p.runs:
            p.runs[0].text = value
            for r in p.runs[1:]:
                r.text = ""
        else:
            p.text = value


def _add_textbox(slide, left, top, width, height, text: str, size: int = 12):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "DejaVu Sans"
    run.font.size = Pt(size)
    return shape


def _build_internal_template(path: Path) -> None:
    """Build a publication-safe letter template from source at runtime."""
    prs = Presentation()
    prs.slide_width = Inches(8.27)
    prs.slide_height = Inches(11.69)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_textbox(slide, Inches(0.6), Inches(0.35), Inches(7.0), Inches(0.45), "WORLD v6 / Saeed Farokhi", 10)
    _add_textbox(slide, Inches(5.1), Inches(1.0), Inches(2.4), Inches(0.35), "تاریخ:", 11)
    _add_textbox(slide, Inches(5.1), Inches(1.4), Inches(2.4), Inches(0.35), "شماره:", 11)
    _add_textbox(slide, Inches(0.8), Inches(2.0), Inches(6.7), Inches(0.45), "گیرنده:", 12)
    _add_textbox(slide, Inches(0.8), Inches(2.55), Inches(6.7), Inches(0.45), "موضوع:", 12)

    body = _add_textbox(slide, Inches(0.8), Inches(3.25), Inches(6.7), Inches(5.2), "", 12)
    body_frame = body.text_frame
    body_frame.clear()
    for i in range(8):
        p = body_frame.paragraphs[0] if i == 0 else body_frame.add_paragraph()
        r = p.add_run()
        r.text = ""
        r.font.name = "DejaVu Sans"
        r.font.size = Pt(12)

    signature = _add_textbox(slide, Inches(4.7), Inches(9.0), Inches(2.8), Inches(1.2), "", 11)
    sig_frame = signature.text_frame
    sig_frame.clear()
    for i in range(3):
        p = sig_frame.paragraphs[0] if i == 0 else sig_frame.add_paragraph()
        r = p.add_run()
        r.text = ""
        r.font.name = "DejaVu Sans"
        r.font.size = Pt(11)

    prs.save(path)


def render_letter_pdf(req, template_pptx: str | Path, output_pdf: str | Path) -> Path:
    template_pptx, output_pdf = Path(template_pptx), Path(output_pdf)
    output_pdf.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="secretary_letter_") as td:
        td = Path(td)
        work = td / "letter.pptx"
        if template_pptx.is_file():
            shutil.copy2(template_pptx, work)
        else:
            _build_internal_template(work)

        prs = Presentation(work)
        slide = prs.slides[0]
        if len(slide.shapes) < 7:
            raise RuntimeError("letter template must expose at least seven ordered shapes")
        _replace_single_run(slide.shapes[1], f"تاریخ: {req.document_date}")
        _replace_single_run(slide.shapes[2], f"شماره: \u202a{req.document_no}\u202c")
        _replace_single_run(slide.shapes[3], f"گیرنده: {req.recipient}")
        _replace_single_run(slide.shapes[4], f"موضوع: {req.subject}")
        _replace_body(slide.shapes[5], req.body)
        _replace_signature(slide.shapes[6], req.signatory, req.company)
        prs.save(work)

        office_binary = shutil.which("libreoffice") or shutil.which("soffice")
        if office_binary is None:
            raise RuntimeError("LibreOffice/soffice is required to render letter PDF")
        office_profile = td / "libreoffice-profile"
        subprocess.run([
            office_binary,
            f"-env:UserInstallation={office_profile.as_uri()}",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(td),
            str(work),
        ], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        generated = td / "letter.pdf"
        if not generated.exists():
            raise RuntimeError("LibreOffice did not generate letter.pdf")
        shutil.copy2(generated, output_pdf)
    return output_pdf
